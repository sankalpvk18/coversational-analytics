from pptx import Presentation
from pptx.util import Inches
import plotly.io as pio
import pandas as pd
import tempfile
import tempfile
import os

# Assuming your project folder is named "MyProject"
project_folder = os.getcwd()


def create_pptx(slides_data):

    # Create a presentation object
    prs = Presentation()

    for slide_data in slides_data:
        question, dataframe, insight, chart = slide_data

        # Add a slide
        slide_layout = prs.slide_layouts[5]  # Use a layout with a title and content
        slide = prs.slides.add_slide(slide_layout)

        # Add the question as the slide title
        title = slide.shapes.title
        title.text = question

        # Add the insight as the content
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(8.5), Inches(1)
        )
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = insight

        # Add the dataframe as a table
        rows, cols = dataframe.shape
        table = slide.shapes.add_table(
            rows + 1, cols, Inches(1), Inches(3), Inches(8.5), Inches(2.5)
        ).table

        # Add the column names
        for col_idx, col_name in enumerate(dataframe.columns):
            table.cell(0, col_idx).text = col_name

        # Add the data
        for row_idx in range(rows):
            for col_idx in range(cols):
                table.cell(row_idx + 1, col_idx).text = str(
                    dataframe.iat[row_idx, col_idx]
                )

        # Save the Plotly chart as an image and add it to the slide
        if chart is not None:
            import os

            # Create a temporary file inside the project folder
            with tempfile.NamedTemporaryFile(
                suffix=".png", dir=project_folder, delete=False
            ) as tmpfile:
                chart_image_path = tmpfile.name
                pio.write_image(chart, chart_image_path)
                slide.shapes.add_picture(
                    chart_image_path, Inches(1), Inches(5.5), Inches(8.5), Inches(3)
                )

    # Save the presentation
    ppt_path = "presentation.pptx"

    prs.save(ppt_path)
    return ppt_path
